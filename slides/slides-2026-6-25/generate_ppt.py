#!/usr/bin/env python3
"""Generate a Nature/Science-style minimal academic PPT for MACR (Chinese)."""
from __future__ import annotations

import json
from pathlib import Path

import matplotlib
import matplotlib.colors
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.slide import Slide
from pptx.util import Inches, Pt

matplotlib.rcParams["font.family"] = ["Noto Sans CJK SC", "sans-serif"]
matplotlib.rcParams["axes.unicode_minus"] = False

ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "assets"
ASSETS.mkdir(exist_ok=True)

PRJ_ROOT = ROOT.parent.parent
SUMMARY_FILE = PRJ_ROOT / "outputs" / "quixbugs" / "ablation_summary.json"
OUT_PPTX = ROOT / "MACR_中期答辩_2026_06_25.pptx"

# --- Palette: Nature/Science minimal ---
PRUSSIAN = RGBColor(26, 54, 93)       # #1A365D
BG = RGBColor(250, 250, 250)          # #FAFAFA
PANEL = RGBColor(255, 255, 255)
TEXT = RGBColor(33, 33, 33)           # #212121
MUTED = RGBColor(120, 120, 120)
LINE = RGBColor(200, 200, 200)
RUST = RGBColor(183, 65, 14)          # #B7410E
GOLD = RGBColor(212, 175, 55)         # #D4AF37
LIGHT_BLUE = RGBColor(230, 240, 250)

FONT = "Microsoft YaHei"
FOOTER = "MACR · 硕士论文中期答辩 · 2026.06"


def load_summary() -> list[dict]:
    return json.loads(SUMMARY_FILE.read_text())


# ---------------------------------------------------------------------------
# Charts: flat, vector-friendly
# ---------------------------------------------------------------------------
def make_success_chart(summary: list[dict]) -> Path:
    configs = [c["configuration"] for c in summary]
    plausible = [c["plausible"] for c in summary]
    generated = [c["patch_generated"] for c in summary]

    x = np.arange(len(configs))
    width = 0.35

    fig, ax = plt.subplots(figsize=(10, 5))
    fig.patch.set_facecolor("#FAFAFA")
    ax.set_facecolor("#FAFAFA")
    bars1 = ax.bar(x - width / 2, plausible, width, label="Plausible", color="#1A365D", edgecolor="white", linewidth=0.5)
    bars2 = ax.bar(x + width / 2, generated, width, label="Generated", color="#D4AF37", edgecolor="white", linewidth=0.5)
    ax.set_ylabel("Bug 数量", fontsize=12, color="#212121")
    ax.set_title("QuixBugs 消融实验结果（n=40）", fontsize=15, fontweight="bold", color="#1A365D", pad=12)
    ax.set_xticks(x)
    ax.set_xticklabels(configs, rotation=15, ha="right", fontsize=10, color="#212121")
    ax.tick_params(axis="y", colors="#212121")
    ax.set_ylim(0, 42)
    ax.legend(frameon=False, fontsize=11)
    ax.axhline(40, color="#CCCCCC", linestyle="--", linewidth=0.8)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#CCCCCC")
    ax.spines["bottom"].set_color("#CCCCCC")
    ax.bar_label(bars1, padding=3, fontsize=10, color="#1A365D")
    ax.bar_label(bars2, padding=3, fontsize=10, color="#8A6D1B")
    plt.tight_layout()
    out = ASSETS / "chart_success_rates.png"
    fig.savefig(out, dpi=220, facecolor="#FAFAFA")
    plt.close(fig)
    return out


def make_failure_matrix(summary: list[dict]) -> Path:
    configs = [c["configuration"] for c in summary]
    bugs = [b["bug_id"] for b in summary[0]["by_bug"]]
    matrix = np.zeros((len(configs), len(bugs)), dtype=int)
    for i, cfg in enumerate(summary):
        for j, bug in enumerate(cfg["by_bug"]):
            if bug["passed"]:
                matrix[i, j] = 2
            elif bug["patch_generated"]:
                matrix[i, j] = 1
            else:
                matrix[i, j] = 0

    fig, ax = plt.subplots(figsize=(16, 5))
    fig.patch.set_facecolor("#FAFAFA")
    ax.set_facecolor("#FAFAFA")
    cmap = matplotlib.colors.ListedColormap(["#B7410E", "#D4AF37", "#1A365D"])
    im = ax.imshow(matrix, cmap=cmap, aspect="auto", vmin=0, vmax=2)
    ax.set_xticks(np.arange(len(bugs)))
    ax.set_xticklabels(bugs, rotation=90, fontsize=7, color="#212121")
    ax.set_yticks(np.arange(len(configs)))
    ax.set_yticklabels(configs, fontsize=9, color="#212121")
    ax.set_title("跨配置失败矩阵", fontsize=15, fontweight="bold", color="#1A365D", pad=12)
    cbar = fig.colorbar(im, ax=ax, ticks=[0, 1, 2], shrink=0.7)
    cbar.ax.set_yticklabels(["未生成", "生成未通过", "通过"], fontsize=9)
    plt.tight_layout()
    out = ASSETS / "chart_failure_matrix.png"
    fig.savefig(out, dpi=220, facecolor="#FAFAFA")
    plt.close(fig)
    return out


# ---------------------------------------------------------------------------
# PPT helpers
# ---------------------------------------------------------------------------
def set_font(run, name=FONT, size=18, bold=False, color=None):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    # Ensure East Asian (Chinese) text uses the same sans-serif font
    rPr = run.font._rPr
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = OxmlElement(tag)
            rPr.append(el)
        el.set("typeface", name)


def slide_bg(slide: Slide, prs: Presentation):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)


def add_footer(slide: Slide, prs: Presentation, number: int):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), prs.slide_height - Inches(0.55), Inches(11.9), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.7), prs.slide_height - Inches(0.45), Inches(11.9), Inches(0.25))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{FOOTER} · {number}"
    set_font(r, FONT, 10, color=MUTED)


def add_top_bar(slide: Slide, title: str, number: int, prs: Presentation):
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(11.5), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_font(r, FONT, 34, True, PRUSSIAN)

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.15), Inches(1.4), Inches(0.04))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RUST
    accent.line.fill.background()

    num_box = slide.shapes.add_textbox(prs.slide_width - Inches(1.1), Inches(0.55), Inches(0.7), Inches(0.5))
    p = num_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = str(number)
    set_font(r, FONT, 18, True, MUTED)

    add_footer(slide, prs, number)


def add_bullets(slide: Slide, items, left=None, top=None, width=None, height=None):
    if left is None:
        left = Inches(0.7)
    if top is None:
        top = Inches(1.55)
    if width is None:
        width = Inches(12)
    if height is None:
        height = Inches(5.4)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, tuple):
            head, subs = item
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.level = 0
            p.space_before = Pt(8)
            p.space_after = Pt(5)
            r = p.add_run()
            r.text = f"■ {head}"
            set_font(r, FONT, 21, True, PRUSSIAN)
            for sub in subs:
                p = tf.add_paragraph()
                p.level = 1
                p.space_after = Pt(3)
                r = p.add_run()
                r.text = f"— {sub}"
                set_font(r, FONT, 19, color=TEXT)
        else:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.level = 0
            p.space_after = Pt(10)
            r = p.add_run()
            r.text = f"■ {item}"
            set_font(r, FONT, 21, color=TEXT)


# ---------------------------------------------------------------------------
# Diagram helpers
# ---------------------------------------------------------------------------
def diagram_box(slide: Slide, text: str, x, y, w, h, fill=PANEL, font_size=14, bold=True, color=PRUSSIAN):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = PRUSSIAN
    shape.line.width = Pt(1.25)
    shape.adjustments[0] = 0.08
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    set_font(r, FONT, font_size, bold, color)
    return shape


def arrow(slide: Slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = PRUSSIAN
    line.line.width = Pt(1)
    line.line.end_arrowhead_width = 2
    line.line.end_arrowhead_length = 2
    return line


def label(slide: Slide, text: str, x, y, w, h, size=12, color=MUTED):
    tb = slide.shapes.add_textbox(x, y, w, h)
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    set_font(r, FONT, size, False, color)


# ---------------------------------------------------------------------------
# Slide constructors
# ---------------------------------------------------------------------------
def add_title_slide(prs, title, subtitle, info):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRUSSIAN
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(12), Inches(1.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    set_font(r, FONT, 46, True, BG)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(12), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = subtitle
    set_font(r, FONT, 28, color=GOLD)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(5.4), Inches(12), Inches(1.2))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = info
    set_font(r, FONT, 18, color=RGBColor(200, 210, 220))

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.7), Inches(6.6), Inches(2.0), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()


def add_divider(prs, number, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRUSSIAN
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(1.6), Inches(1.3))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = number
    set_font(r, FONT, 72, True, BG)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(12), Inches(1.4))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_font(r, FONT, 48, True, BG)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.0), Inches(1.8), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()


def add_content_slide(prs, title, items, number, image_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, prs)
    add_top_bar(slide, title, number, prs)

    if image_path is None:
        add_bullets(slide, items)
    else:
        add_bullets(slide, items, width=Inches(6.7))
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.6), Inches(1.6), Inches(5.0), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = PANEL
        card.line.color.rgb = LINE
        card.line.width = Pt(1)
        card.adjustments[0] = 0.04
        slide.shapes._spTree.remove(card._element)
        slide.shapes._spTree.insert(3, card._element)
        slide.shapes.add_picture(str(image_path), Inches(7.75), Inches(1.75), width=Inches(4.7))


def add_table_slide(prs, summary, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, prs)
    add_top_bar(slide, "实验结果汇总", number, prs)

    rows, cols = 6, 5
    left = Inches(0.7)
    top = Inches(1.65)
    width = Inches(11.9)
    height = Inches(4.8)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    headers = ["配置", "Plausible", "Generated", "Avg Steps", "Tokens"]
    tokens = {
        "Baseline": 176_408,
        "MACR": 171_464,
        "MACR + Sub-Agent": 312_405,
        "MACR + KG": 196_233,
        "MACR + Sub-Agent + KG": 285_216,
    }
    for i, h in enumerate(headers):
        c = table.cell(0, i)
        c.text = h
        for p in c.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                set_font(r, FONT, 17, True, BG)
        c.fill.solid()
        c.fill.fore_color.rgb = PRUSSIAN

    for ri, cfg in enumerate(summary, start=1):
        name = cfg["configuration"]
        vals = [
            name,
            f"{cfg['plausible']}/{cfg['total']} ({cfg['plausible_rate']*100:.1f}%)",
            f"{cfg['patch_generated']}/{cfg['total']} ({cfg['patch_generated_rate']*100:.1f}%)",
            f"{cfg['average_attempts']:.2f}",
            f"{tokens[name]:,}",
        ]
        for ci, v in enumerate(vals):
            c = table.cell(ri, ci)
            c.text = v
            for p in c.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for r in p.runs:
                    set_font(r, FONT, 16, color=TEXT)
            c.fill.solid()
            c.fill.fore_color.rgb = PANEL if ri % 2 else RGBColor(245, 247, 250)


def add_image_slide(prs, title, image_path, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, prs)
    add_top_bar(slide, title, number, prs)

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.55), Inches(11.5), Inches(5.4))
    card.fill.solid()
    card.fill.fore_color.rgb = PANEL
    card.line.color.rgb = LINE
    card.line.width = Pt(1)
    card.adjustments[0] = 0.03
    slide.shapes._spTree.remove(card._element)
    slide.shapes._spTree.insert(3, card._element)

    slide.shapes.add_picture(str(image_path), Inches(1.05), Inches(1.7), width=Inches(11.2))


def add_architecture_slide(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, prs)
    add_top_bar(slide, "MACR 总体架构", number, prs)

    left_col = Inches(1.0)
    mid_x = Inches(5.8)
    right_col = Inches(9.8)
    top_y = Inches(2.0)
    mid_y = Inches(3.5)
    bot_y = Inches(5.3)
    bw = Inches(2.6)
    bh = Inches(0.75)

    diagram_box(slide, "ReAct\nCoordinator", mid_x, top_y, bw, bh, LIGHT_BLUE)
    diagram_box(slide, "Tool Registry\n(direct_patch / run_tests / query_kg)", mid_x, mid_y, bw, bh)
    diagram_box(slide, "Knowledge\nGraph", left_col, mid_y, bw, bh)
    diagram_box(slide, "Level-2\nSub-Agent", right_col, mid_y, bw, bh)
    diagram_box(slide, "Trace\nLogger", left_col, bot_y, bw, bh)
    diagram_box(slide, "QuixBugs\nDataset", right_col, bot_y, bw, bh)
    diagram_box(slide, "LLM Backend\nDeepSeek-V4", mid_x, bot_y, bw, bh)

    arrow(slide, mid_x + bw / 2, top_y + bh, mid_x + bw / 2, mid_y)
    arrow(slide, mid_x, mid_y + bh / 2, left_col + bw, mid_y + bh / 2)
    arrow(slide, mid_x + bw, mid_y + bh / 2, right_col, mid_y + bh / 2)
    arrow(slide, mid_x + bw / 2, mid_y + bh, mid_x + bw / 2, bot_y)
    arrow(slide, left_col + bw / 2, mid_y + bh, left_col + bw / 2, bot_y)
    arrow(slide, right_col + bw / 2, mid_y + bh, right_col + bw / 2, bot_y)


def add_workflow_slide(prs, number):
    """MACR 工作流程图（横向流程）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, prs)
    add_top_bar(slide, "MACR 工作流程", number, prs)

    y = Inches(3.0)
    h = Inches(0.9)
    w = Inches(2.0)
    gap = Inches(2.3)
    x0 = Inches(0.7)

    stages = [
        ("Bug\nSample", PANEL),
        ("KG\nBuilder", PANEL),
        ("ReAct\nCoordinator", LIGHT_BLUE),
        ("Patch\nGeneration", PANEL),
        ("Test\nValidation", PANEL),
        ("Trace\nLogger", PANEL),
    ]
    xs = []
    for i, (text, fill) in enumerate(stages):
        x = x0 + i * gap
        xs.append(x)
        diagram_box(slide, text, x, y, w, h, fill, font_size=13)
        if i < len(stages) - 1:
            arrow(slide, x + w, y + h / 2, x + gap, y + h / 2)

    # feedback loop from Test to Coordinator
    feedback_y = Inches(4.4)
    label(slide, "失败反馈 / 迭代", Inches(4.5), Inches(4.5), Inches(3.5), Inches(0.3), size=12, color=RUST)
    arrow(slide, xs[4] + w / 2, y + h, xs[4] + w / 2, feedback_y)
    arrow(slide, xs[4] + w / 2, feedback_y, xs[2] + w / 2, feedback_y)
    arrow(slide, xs[2] + w / 2, feedback_y, xs[2] + w / 2, y + h)


def add_react_loop_slide(prs, number):
    """ReAct 循环图"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, prs)
    add_top_bar(slide, "ReAct 决策循环", number, prs)

    cx = Inches(6.7)
    cy = Inches(3.75)
    radius = Inches(2.0)
    w = Inches(2.0)
    h = Inches(0.7)

    positions = [
        ("Observation", cx, cy - radius),
        ("Thought", cx + radius * 0.95, cy - radius * 0.3),
        ("Action", cx + radius * 0.75, cy + radius * 0.75),
        ("Tool / Patch", cx - radius * 0.75, cy + radius * 0.75),
        ("Feedback", cx - radius * 0.95, cy - radius * 0.3),
    ]
    for text, x, y in positions:
        # center the box at (x,y)
        diagram_box(slide, text, x - w / 2, y - h / 2, w, h, LIGHT_BLUE if text == "Thought" else PANEL, font_size=14)

    # circular arrows
    for i in range(len(positions)):
        x1, y1 = positions[i][1], positions[i][2]
        x2, y2 = positions[(i + 1) % len(positions)][1], positions[(i + 1) % len(positions)][2]
        # shorten to box edge
        dx, dy = x2 - x1, y2 - y1
        length = (dx**2 + dy**2) ** 0.5
        if length == 0:
            continue
        offset = Inches(1.1)
        sx = x1 + dx / length * offset
        sy = y1 + dy / length * offset
        ex = x2 - dx / length * offset
        ey = y2 - dy / length * offset
        arrow(slide, sx, sy, ex, ey)


def add_kg_pipeline_slide(prs, number):
    """KG 构建流程"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, prs)
    add_top_bar(slide, "Knowledge Graph 构建流程", number, prs)

    y = Inches(2.2)
    h = Inches(0.85)
    w = Inches(2.4)
    gap = Inches(2.7)
    x0 = Inches(0.7)

    stages = [
        ("Python\n源文件", PANEL),
        ("AST\n解析", PANEL),
        ("实体抽取\n函数 / 类 / 变量", PANEL),
        ("关系抽取\n调用 / 引用", PANEL),
        ("图构建\nG = (V, E)", LIGHT_BLUE),
    ]
    for i, (text, fill) in enumerate(stages):
        x = x0 + i * gap
        diagram_box(slide, text, x, y, w, h, fill, font_size=13)
        if i < len(stages) - 1:
            arrow(slide, x + w, y + h / 2, x + gap, y + h / 2)

    # query return
    query_y = Inches(4.0)
    diagram_box(slide, "query_kg(q) → 相关实体子集", Inches(4.2), query_y, Inches(5.0), Inches(0.7), PANEL, font_size=14)
    arrow(slide, x0 + 4 * gap + w / 2, y + h, x0 + 4 * gap + w / 2, query_y)
    label(slide, "Coordinator 查询", Inches(7.5), Inches(3.55), Inches(2.5), Inches(0.3), size=12, color=RUST)


def add_subagent_flow_slide(prs, number):
    """Sub-Agent 调用流程"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, prs)
    add_top_bar(slide, "Sub-Agent 调用机制", number, prs)

    # top: coordinator decision
    diagram_box(slide, "Coordinator\n判断是否调用 Sub-Agent", Inches(4.5), Inches(1.8), Inches(4.3), Inches(0.8), LIGHT_BLUE, font_size=14)

    y = Inches(3.2)
    h = Inches(1.0)
    w = Inches(2.6)
    gap = Inches(3.0)
    x0 = Inches(1.0)

    agents = ["加载配置\nMarkdown + YAML", "实例化\n隔离上下文", "执行\nmini ReAct", "返回\n结构化摘要"]
    for i, text in enumerate(agents):
        x = x0 + i * gap
        diagram_box(slide, text, x, y, w, h, PANEL, font_size=13)
        if i == 0:
            arrow(slide, Inches(6.65), Inches(2.6), x + w / 2, y)
        if i < len(agents) - 1:
            arrow(slide, x + w, y + h / 2, x + gap, y + h / 2)

    # return to coordinator
    arrow(slide, x0 + 3 * gap + w / 2, y + h, Inches(6.65), Inches(5.0))
    diagram_box(slide, "Coordinator 融合摘要\n继续主循环", Inches(4.5), Inches(5.0), Inches(4.3), Inches(0.8), LIGHT_BLUE, font_size=14)


def add_formula_slide(prs, number):
    """关键公式"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, prs)
    add_top_bar(slide, "形式化定义", number, prs)

    formulas = [
        ("ReAct 动作选择", "a_t = argmaxₐ P(a | s_t, h_{t-1})"),
        ("Patch 生成目标", "p* = argmaxₚ∈𝒫  P(pass(T) | p, c, q)"),
        ("Knowledge Graph", "G = (V, E),  V = Func ∪ Class ∪ Var,  E = Call ∪ Ref"),
        ("KG 查询", "R(q) = { e ∈ V | relevance(e, q) ≥ θ }"),
        ("Sub-Agent 输出", "s_i = SubAgent(config_i, ctx, instr)"),
    ]

    y0 = Inches(1.7)
    h = Inches(0.75)
    gap = Inches(0.95)
    for i, (name, formula) in enumerate(formulas):
        y = y0 + i * gap
        # name tag
        tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), y, Inches(2.4), h)
        tag.fill.solid()
        tag.fill.fore_color.rgb = PRUSSIAN
        tag.line.fill.background()
        tag.adjustments[0] = 0.1
        tf = tag.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = name
        set_font(r, FONT, 14, True, BG)

        # formula card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.3), y, Inches(9.3), h)
        card.fill.solid()
        card.fill.fore_color.rgb = PANEL
        card.line.color.rgb = LINE
        card.line.width = Pt(1)
        card.adjustments[0] = 0.05
        tf = card.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = formula
        set_font(r, FONT, 18, True, PRUSSIAN)


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------
def build_ppt():
    summary = load_summary()
    chart = make_success_chart(summary)
    make_failure_matrix(summary)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides_data = []

    def add(fn, *args):
        slides_data.append((fn, args))

    add(add_title_slide, prs, "基于多智能体协作的自动化程序修复研究", "硕士论文中期答辩", "汇报人：XXX\n导师：XXX 教授\n2026年6月25日")
    add(add_content_slide, prs, "汇报提纲", [
        "1. 论文概况（背景、目标）",
        "2. 研究进展概览",
        "3. 已完成的核心工作",
        "4. 存在的问题与困惑",
        "5. 后续研究计划",
        "6. 总结与致谢",
    ], 2)

    add(add_divider, prs, "01", "论文概况")
    add(add_content_slide, prs, "研究背景与目标", [
        ("研究背景", [
            "软件缺陷修复成本高，传统 APR 泛化能力有限",
            "大语言模型展现代码理解与生成潜力",
            "单一 LLM 代理存在上下文与推理局限",
            "多智能体协作可分工完成定位、生成、验证、评审",
        ]),
        ("研究目标", [
            "构建 MACR（Multi-Agent Collaborative Repair）框架",
            "引入 Knowledge Graph 与 Level-2 Sub-Agent 增强机制",
            "在 QuixBugs 数据集上验证有效性",
        ]),
    ], 4)

    add(add_divider, prs, "02", "研究进展概览")
    add(add_content_slide, prs, "进度对照与完成情况", [
        ("已完成", [
            "文献综述与相关工作梳理",
            "MACR 总体架构设计与模块实现",
            "Coordinator / KG / Sub-Agent / Trace 核心模块",
            "QuixBugs 数据集消融实验与阶段报告",
        ]),
        ("进行中", [
            "失败 case 的 trace 定位与 prompt 优化",
            "Sub-Agent 触发策略的成本-收益调优",
        ]),
        ("待完成", [
            "Defects4J 等更大规模数据集验证",
            "论文撰写、导师审阅与修改",
        ]),
    ], 6)

    add(add_divider, prs, "03", "已完成的核心工作")
    add(add_architecture_slide, prs, 8)
    add(add_workflow_slide, prs, 9)
    add(add_formula_slide, prs, 10)
    add(add_react_loop_slide, prs, 11)
    add(add_content_slide, prs, "ReAct Coordinator 与工具链", [
        ("动作选择", [
            "LLM 输出 ReActAction：tool_name + tool_input + reasoning",
            "支持直接生成 patch 与多步测试-迭代",
        ]),
        ("测试反馈", [
            "run_tests 返回 pytest 详细输出（stdout/stderr）",
            "失败信息驱动下一轮动作选择",
        ]),
        ("执行边界", [
            "max_steps = 5，避免无限循环",
            "测试通过可提前终止",
        ]),
    ], 12)

    add(add_kg_pipeline_slide, prs, 13)
    add(add_content_slide, prs, "Knowledge Graph 增强", [
        ("构建方式", [
            "解析 bug 目录全部 .py 文件",
            "抽取函数、类、变量及调用/引用关系",
        ]),
        ("查询工具 query_kg", [
            "Coordinator 主动查询相关函数与依赖",
            "将图结构信息注入 LLM 上下文",
        ]),
        ("实验效果", [
            "MACR+KG 达到 95% plausible、100% generated",
            "显著优于纯 MACR（80%）与 Baseline（87.5%）",
        ]),
    ], 14)

    add(add_subagent_flow_slide, prs, 15)
    add(add_content_slide, prs, "Level-2 Sub-Agent 机制", [
        ("配置驱动", [
            "Markdown + YAML frontmatter 定义 agent",
            "支持独立 model、tools、max_steps 配置",
        ]),
        ("内置角色", [
            "localize、generate_patch、validate_patch",
            "explore、review、security_audit、planner",
        ]),
        ("当前问题", [
            "explore 调用过于频繁，token 开销大",
            "与 KG 组合未产生协同增益",
        ]),
    ], 16)

    add(add_content_slide, prs, "Trace 与实验设计", [
        ("Trace 记录", [
            "run_start/end、react_step、llm_request/response、subagent_spawn",
        ]),
        ("实验脚本", [
            "run_ablation.py：5 配置消融",
            "analyze_traces.py：成功率、token、调用统计",
        ]),
        ("实验设置", [
            "数据集：QuixBugs Python（40 bugs）",
            "模型：SiliconFlow DeepSeek-V4-Flash",
            "指标：Plausible / Generated / Avg Steps / Tokens",
        ]),
    ], 17)

    add(add_table_slide, prs, summary, 18)
    add(add_image_slide, prs, "实验结果可视化", chart, 19)

    add(add_content_slide, prs, "关键发现", [
        ("Knowledge Graph 提升最明显", [
            "MACR+KG：38/40 plausible，40/40 generated",
            "比 Baseline 多 3 个 plausible，且全部生成 patch",
        ]),
        ("Sub-Agent 成本偏高、收益有限", [
            "token 消耗 312k，约为 MACR+KG 的 1.6 倍",
            "plausible 率 90%，未超过 KG 单独使用",
        ]),
        ("纯 MACR 低于 Baseline", [
            "说明 ReAct loop 与 system prompt 仍有优化空间",
        ]),
        ("硬 bug", [
            "max_sublist_sum：KG 增强后仍失败",
            "find_first_in_sorted：MACR/Sub-Agent 无法生成",
            "reverse_linked_list / shortest_path_lengths / to_base 反复失败",
        ]),
    ], 20)

    add(add_divider, prs, "04", "存在的问题与困惑")
    add(add_content_slide, prs, "技术难点与困惑", [
        ("Sub-Agent 触发策略", [
            "何时调用 explore / generate_patch 尚不明确",
            "频繁调用导致 token 成本显著上升",
        ]),
        ("Knowledge Graph 粒度", [
            "部分硬 bug 在 KG 增强后仍失败",
            "图信息是否过粗或存在误导？",
        ]),
        ("ReAct 稳定性", [
            "偶发 JSON 解析失败、测试超时",
            "prompt 与容错机制需加强",
        ]),
        ("数据集规模", [
            "仅验证 QuixBugs，缺乏 Defects4J 等真实大型项目验证",
        ]),
    ], 22)

    add(add_divider, prs, "05", "后续研究计划")
    add(add_content_slide, prs, "后续时间表与拟解决方案", [
        ("2026.06", [
            "失败 case trace 分析",
            "prompt 与 action schema 优化",
        ]),
        ("2026.07", [
            "重设计 Sub-Agent 触发条件（仅在失败/定位时调用）",
            "引入代价模型避免无效探索",
            "针对性消融实验",
        ]),
        ("2026.08", [
            "Defects4J 扩展实验",
            "成本-效益分析与对比实验",
        ]),
        ("2026.09", [
            "论文撰写、导师审阅与修改",
            "准备答辩",
        ]),
    ], 24)

    add(add_content_slide, prs, "预期产出", [
        ("学术论文", [
            "一篇关于 MACR 的硕士学位论文",
            "拟投稿 1 篇软件工程领域会议/期刊论文",
        ]),
        ("系统贡献", [
            "开源 MACR 框架与实验复现脚本",
            "完整的 trace 数据集与分析工具",
        ]),
        ("核心结论", [
            "验证 KG 在 LLM-driven APR 中的有效性",
            "明确 Sub-Agent 的适用场景与优化方向",
        ]),
    ], 25)

    add(add_divider, prs, "06", "总结与致谢")
    add(add_content_slide, prs, "总结", [
        "已完成 MACR 框架设计与实现，并通过 QuixBugs 消融实验验证",
        "Knowledge Graph 是提升修复成功率最有效的模块",
        "Sub-Agent 当前成本高、收益有限，需重新设计触发策略",
        "后续将聚焦 hard case 分析、Sub-Agent 优化与 Defects4J 扩展实验",
        "计划 2026 年 9 月完成论文撰写与答辩准备",
    ], 27)

    add(add_content_slide, prs, "致谢", [
        "感谢导师的悉心指导与建议",
        "感谢课题组同学的帮助与讨论",
        "恳请各位老师批评指正！",
    ], 28)

    for fn, args in slides_data:
        fn(*args)

    prs.save(OUT_PPTX)
    print(f"PPT saved: {OUT_PPTX} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build_ppt()
